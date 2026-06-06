"""Content extraction across all supported attachment types.

Reuses the proven ``app.services.gpt_service.extract_text_from_file`` for the
types it already handles (PDF, DOC/DOCX, TXT, CSV/XLS/XLSX, images→base64,
audio/video→transcript) and ADDS the rest the assistant needs:

  * source code / markup (md, rtf, json, xml, yaml, js, ts, py, ...) → fenced text
  * presentations (ppt/pptx) → slide text via python-pptx
  * archives (zip) → safe walk: folder tree + extracted text of contained files

Returns a uniform dict: {"text": str, "images": list[data-uri], "metadata": dict}.
Never raises — extraction failures are captured as a note so the upload still
completes with a sensible status.
"""
from __future__ import annotations

import os
import zipfile
from typing import Dict, List

from app.assistant.config import config

# Extension → coarse kind, used for UI icons and routing.
_KIND_MAP = {
    "document": {"pdf", "doc", "docx", "txt", "md", "rtf"},
    "spreadsheet": {"xls", "xlsx", "csv"},
    "presentation": {"ppt", "pptx"},
    "image": {"jpg", "jpeg", "png", "webp", "gif"},
    "audio": {"mp3", "wav", "aac", "ogg", "m4a", "flac"},
    "video": {"mp4", "mov", "avi", "mkv", "webm"},
    "code": {
        "js", "jsx", "ts", "tsx", "py", "java", "php", "cpp", "c", "h", "cs",
        "go", "rb", "rs", "kt", "swift", "sql", "sh", "json", "xml", "yaml",
        "yml", "html", "css",
    },
    "archive": {"zip", "rar", "7z"},
}

# Plain-text-ish extensions we read directly as UTF-8 (code + lightweight docs).
_TEXT_EXTS = _KIND_MAP["code"] | {"txt", "md"}

# Language hint for code fences (cosmetic; helps the model).
_LANG = {
    "py": "python", "js": "javascript", "jsx": "jsx", "ts": "typescript",
    "tsx": "tsx", "java": "java", "php": "php", "cpp": "cpp", "c": "c",
    "cs": "csharp", "go": "go", "rb": "ruby", "rs": "rust", "kt": "kotlin",
    "swift": "swift", "sql": "sql", "sh": "bash", "json": "json", "xml": "xml",
    "yaml": "yaml", "yml": "yaml", "html": "html", "css": "css",
}


def ext_of(filename: str) -> str:
    return (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()


def kind_of(filename: str) -> str:
    e = ext_of(filename)
    for kind, exts in _KIND_MAP.items():
        if e in exts:
            return kind
    return "other"


def _read_text_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _extract_pptx(path: str) -> str:
    """Slide-by-slide text from a presentation."""
    try:
        from pptx import Presentation
    except Exception:  # python-pptx not installed
        return "[Presentation: install python-pptx to extract slide text]"
    prs = Presentation(path)
    out: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        lines.append(txt)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    lines.append(" | ".join(cells))
        out.append("\n".join(lines))
    return "\n\n".join(out)


def _safe_zip_members(zf: zipfile.ZipFile) -> List[zipfile.ZipInfo]:
    """Return safe members, enforcing entry-count, total-size and zip-slip guards."""
    members: List[zipfile.ZipInfo] = []
    total = 0
    for info in zf.infolist():
        if info.is_dir():
            continue
        name = info.filename
        # zip-slip / absolute path guard.
        if name.startswith("/") or ".." in name.replace("\\", "/").split("/"):
            continue
        total += info.file_size
        if total > config.ZIP_MAX_TOTAL_BYTES:
            break
        members.append(info)
        if len(members) >= config.ZIP_MAX_ENTRIES:
            break
    return members


async def _extract_zip(path: str, filename: str) -> Dict:
    """Walk an archive safely: list its tree and extract text from contained
    supported files (one level of recursion — nested archives are listed only)."""
    text_parts: List[str] = []
    images: List[str] = []
    tree: List[str] = []
    try:
        with zipfile.ZipFile(path) as zf:
            members = _safe_zip_members(zf)
            for info in zf.infolist():
                if not info.is_dir():
                    tree.append(info.filename)
            text_parts.append("[Archive structure]\n" + "\n".join(sorted(tree)[:200]))

            import tempfile
            for info in members:
                inner_ext = ext_of(info.filename)
                if inner_ext in config.BLOCKED_EXTENSIONS or inner_ext not in config.ALLOWED_EXTENSIONS:
                    continue
                if inner_ext in {"zip", "rar", "7z"}:
                    continue  # don't recurse into nested archives
                suffix = f".{inner_ext}" if inner_ext else ""
                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                    tmp.write(zf.read(info))
                    tmp_path = tmp.name
                try:
                    sub = await extract(tmp_path, os.path.basename(info.filename))
                    if sub["text"].strip():
                        text_parts.append(f"\n# {info.filename}\n{sub['text']}")
                    images.extend(sub.get("images") or [])
                finally:
                    try:
                        os.remove(tmp_path)
                    except OSError:
                        pass
    except zipfile.BadZipFile:
        text_parts.append(f"[Archive: {filename} could not be read]")
    return {
        "text": "\n".join(text_parts),
        "images": images[: config.MAX_IMAGES_PER_TURN],
        "metadata": {"archive_entries": len(tree)},
    }


async def extract(local_path: str, filename: str) -> Dict:
    """Extract {text, images, metadata} from a file. Never raises."""
    ext = ext_of(filename)
    metadata: Dict = {"kind": kind_of(filename), "ext": ext}

    try:
        # Archives — safe walk.
        if ext == "zip":
            result = await _extract_zip(local_path, filename)
            result["metadata"] = {**metadata, **result.get("metadata", {})}
            return result
        if ext in {"rar", "7z"}:
            return {
                "text": f"[Archive: {filename} — {ext.upper()} extraction not yet supported; "
                        f"upload a .zip to analyse contents]",
                "images": [],
                "metadata": metadata,
            }

        # Presentations.
        if ext in {"ppt", "pptx"}:
            return {"text": _extract_pptx(local_path), "images": [], "metadata": metadata}

        # Code / markup / plain text → read natively (no heavy deps). Code gets a
        # language-fenced block; txt/md stay plain. Handling these here also keeps
        # simple text files working even if optional media libs aren't installed.
        if ext in _TEXT_EXTS:
            body = _read_text_file(local_path)
            lang = _LANG.get(ext, "")
            text = f"```{lang}\n{body}\n```" if lang else body
            metadata["lines"] = body.count("\n") + 1
            return {"text": text, "images": [], "metadata": metadata}

        # rtf → strip control words to plain text (best-effort, no extra dep).
        if ext == "rtf":
            return {"text": _strip_rtf(_read_text_file(local_path)), "images": [], "metadata": metadata}

        # Audio / video → transcription. Handled explicitly so each failure mode
        # (no ffmpeg, no speech_recognition, silent file) yields a clear message
        # the assistant can relay, instead of an empty or cryptic result.
        if metadata["kind"] in {"audio", "video"}:
            return await _extract_media(local_path, filename, metadata)

        # Everything else (pdf, doc/docx, csv/xls/xlsx, images) is handled by the
        # shared, battle-tested extractor.
        from app.services.gpt_service import extract_text_from_file

        result = await extract_text_from_file(local_path, filename)
        return {
            "text": result.get("text", "") or "",
            "images": result.get("images", []) or [],
            "metadata": metadata,
        }
    except Exception as exc:  # noqa: BLE001 — extraction must never crash the job
        return {
            "text": f"[Could not extract content from {filename}: {exc}]",
            "images": [],
            "metadata": {**metadata, "extraction_error": str(exc)},
        }


async def _extract_media(local_path: str, filename: str, metadata: Dict) -> Dict:
    """Transcribe an audio/video file, with explicit guards for the optional
    dependencies (ffmpeg binary + speech_recognition). Always returns readable
    text describing the outcome so the assistant can explain it to the user."""
    from app.services.media_tools import ffmpeg_available

    kind = metadata.get("kind", "media")

    if not ffmpeg_available():
        return {
            "text": f"[{filename}: this {kind} could not be transcribed because the "
                    f"ffmpeg media tool is not installed on the server.]",
            "images": [],
            "metadata": {**metadata, "transcription": "ffmpeg_missing"},
        }

    try:
        from app.services import transcription_service as ts
    except Exception as exc:  # noqa: BLE001
        return {
            "text": f"[{filename}: transcription is unavailable on the server ({exc}).]",
            "images": [],
            "metadata": {**metadata, "transcription": "unavailable"},
        }

    if getattr(ts, "sr", None) is None:
        return {
            "text": f"[{filename}: this {kind} could not be transcribed because the "
                    f"speech_recognition package is not installed on the server.]",
            "images": [],
            "metadata": {**metadata, "transcription": "sr_missing"},
        }

    transcript = await ts.transcribe_media_file(local_path)
    if transcript and transcript.strip():
        return {
            "text": f"[Transcript of {filename}]\n{transcript.strip()}",
            "images": [],
            "metadata": {**metadata, "transcription": "ok"},
        }
    return {
        "text": f"[{filename}: no speech could be detected in this {kind} file.]",
        "images": [],
        "metadata": {**metadata, "transcription": "empty"},
    }


def _strip_rtf(data: str) -> str:
    """Very small RTF→text fallback (removes control words and braces)."""
    import re

    text = re.sub(r"\\'[0-9a-fA-F]{2}", "", data)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("{", "").replace("}", "")
    return text.strip()
